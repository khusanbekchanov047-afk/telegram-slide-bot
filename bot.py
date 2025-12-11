import asyncio
import os
import textwrap
import random
from aiogram import Bot, Dispatcher, types
from aiogram.filters import Command
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

TOKEN = "8345279215:AAHQaef_NeQer9nKn0T82GQxWp7jhewbv_Y"

bot = Bot(token=TOKEN)
dp = Dispatcher()

user_data = {}

def wrap_text(text, max_chars=50):
    return "\n".join(textwrap.wrap(text, width=max_chars))

def extract_title(text):
    if '.' in text:
        title, rest = text.split('.', 1)
        return title.strip() + '.', rest.strip()
    else:
        return text, ""

# Turli mavzu ranglari
title_colors = [
    RGBColor(0,0,0),
    RGBColor(255,0,0),
    RGBColor(0,128,0),
    RGBColor(0,0,255),
    RGBColor(128,0,128)
]

# Turli fon ranglari
background_colors = [
    RGBColor(255, 255, 255), # oq
    RGBColor(240, 248, 255), # alice blue
    RGBColor(255, 239, 213), # papaya whip
    RGBColor(245, 245, 220), # beige
    RGBColor(224, 255, 255), # light cyan
]

@dp.message(Command("start"))
async def start(message: types.Message):
    user_data[message.from_user.id] = {"texts": [], "images": []}
    await message.answer(
        "Botga matn yuboring.\nRasm yuborsangiz ham qabul qiladi.\nOxirida /finish bosing — pptx slayd tayyorlab beraman."
    )

@dp.message(Command("finish"))
async def finish(message: types.Message):
    uid = message.from_user.id
    data = user_data.get(uid)

    if not data or (not data["texts"] and not data["images"]):
        await message.answer("Hech narsa yubormadingiz.")
        return

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for i, text in enumerate(data["texts"]):
        slide = prs.slides.add_slide(blank_slide_layout)

        # --- Slayd fonini turlicha rang berish ---
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = random.choice(background_colors)

        # --- Bosh mavzu ---
        title_text, body_text = extract_title(text)
        if title_text:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1.2))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = title_text
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = random.choice(title_colors)
            p.alignment = 1  # center

        # --- Matn joylashuvi ---
        # Agar rasm mavjud bo‘lsa → matn o‘ng tomonda
        if i < len(data["images"]):
            img_path = data["images"][i]
            # Rasm chapga
            slide.shapes.add_picture(img_path, Inches(0.5), Inches(2), width=Inches(4))
            # Matn o‘ngga
            txBox2 = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4), Inches(3))
        else:
            # Rasm yo‘q → matn markaz
            txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(3))

        if body_text:
            tf2 = txBox2.text_frame
            p2 = tf2.add_paragraph()
            p2.text = wrap_text(body_text, max_chars=50)
            p2.font.size = Pt(24)
            p2.font.name = 'Arial'

    # --- Fayl saqlash ---
    filename = f"{uid}_presentation.pptx"
    prs.save(filename)

    await message.answer_document(types.FSInputFile(filename))

    for img in data["images"]:
        os.remove(img)
    os.remove(filename)

    user_data[uid] = {"texts": [], "images": []}
    await message.answer("Tayyor! Yana matn yoki rasm yuborishingiz mumkin.")

@dp.message()
async def recv(message: types.Message):
    uid = message.from_user.id
    if uid not in user_data:
        user_data[uid] = {"texts": [], "images": []}

    if message.photo:
        file_id = message.photo[-1].file_id
        file = await bot.get_file(file_id)
        file_path = file.file_path
        downloaded = await bot.download_file(file_path)
        img_name = f"{uid}_{len(user_data[uid]['images'])}.jpg"
        with open(img_name, "wb") as f:
            f.write(downloaded.read())
        user_data[uid]["images"].append(img_name)
        await message.answer("📸 Rasm qabul qilindi")
        return

    user_data[uid]["texts"].append(message.text)
    await message.answer("✏️ Matn qabul qilindi")

async def main():
    await dp.start_polling(bot)

if __name__ == "__main__":
    asyncio.run(main())
